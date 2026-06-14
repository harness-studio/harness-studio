"""Generate Harness-Studio-Overview.pptx — current model (v0.2.0-alpha.1)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── palette ──────────────────────────────────────────────────────────────────
BG       = RGBColor(0x0d, 0x11, 0x17)   # near-black
SURFACE  = RGBColor(0x16, 0x1b, 0x22)   # card background
BLUE     = RGBColor(0x58, 0xa6, 0xff)   # accent blue
ORANGE   = RGBColor(0xf0, 0x88, 0x3e)   # accent orange
GREEN    = RGBColor(0x3f, 0xb9, 0x50)   # accent green
WHITE    = RGBColor(0xff, 0xff, 0xff)
MUTED    = RGBColor(0x8b, 0x94, 0x9e)   # secondary text
SLIDE_W  = Inches(13.33)
SLIDE_H  = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # totally blank


def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def box(slide, left, top, width, height, text="", size=18, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, bg=None, wrap=True):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return txb


def label(slide, left, top, text, size=10, color=BLUE, bold=True):
    box(slide, left, top, 3, 0.35, text, size=size, color=color, bold=bold)


def para(slide, left, top, width, height, text, size=14, color=WHITE):
    box(slide, left, top, width, height, text, size=size, color=color)


def code(slide, left, top, width, height, text, size=11):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = False
    txb.fill.solid()
    txb.fill.fore_color.rgb = SURFACE
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.color.rgb = GREEN
    run.font.name  = "Courier New"
    return txb


def multiline_box(slide, left, top, width, height, lines, default_size=14,
                  default_color=WHITE, bg=None):
    """lines = list of (text, size, color, bold, align)"""
    txb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    if bg:
        txb.fill.solid()
        txb.fill.fore_color.rgb = bg
    first = True
    for item in lines:
        if isinstance(item, str):
            text, sz, col, bld, aln = item, default_size, default_color, False, PP_ALIGN.LEFT
        else:
            text = item[0]
            sz   = item[1] if len(item) > 1 else default_size
            col  = item[2] if len(item) > 2 else default_color
            bld  = item[3] if len(item) > 3 else False
            aln  = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = aln
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(sz)
        run.font.bold  = bld
        run.font.color.rgb = col
    return txb


def hline(slide, top, color=SURFACE):
    ln = slide.shapes.add_shape(1, Inches(0.5), Inches(top), Inches(12.33), Emu(3000))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()


def slide_number(slide, n, total=14):
    box(slide, 12.3, 7.1, 0.9, 0.3, f"{n}/{total}", size=9, color=MUTED,
        align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
box(s, 0.6, 0.4, 12, 0.5, "// AI-FIRST  SOFTWARE  DELIVERY", size=13,
    color=MUTED, bold=False)
box(s, 0.6, 1.1, 11, 1.5, "Harness Studio", size=54, bold=True, color=WHITE)
box(s, 0.6, 2.7, 10, 0.7,
    "A governed, adversarial framework for delivering software with AI agents.",
    size=20, color=MUTED)
hline(s, 3.7)
multiline_box(s, 0.6, 3.9, 11, 1.5, [
    ("Whoever does the work", 28, WHITE, False),
    ('never grades their own “done.”', 28, BLUE, True),
])
box(s, 0.6, 6.8, 7, 0.4,
    "github.com/harness-studio/harness-studio   ·   hssd   ·   v0.2.0-alpha.1",
    size=11, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — THE PROBLEM
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 2)
label(s, 0.6, 0.4, "01 / THE PROBLEM")
box(s, 0.6, 0.8, 11, 0.8, "Two ways AI delivery fails", size=28, bold=True)
hline(s, 1.8)

multiline_box(s, 0.6, 2.0, 5.7, 3.5, [
    ("FAILURE MODE 1", 13, ORANGE, True),
    ("Too fast, no structure", 18, WHITE, True),
    ("", 8, WHITE),
    ('request → agent builds → "done" → ships', 13, MUTED, False),
    ("", 8, WHITE),
    ("Missing: scope validation, architecture,\ntests, adversarial verification.", 13, WHITE),
    ("", 8, WHITE),
    ("Result: bugs in prod, no audit trail,\nno confidence when it counts.", 13, MUTED),
])

multiline_box(s, 6.9, 2.0, 5.8, 3.5, [
    ("FAILURE MODE 2", 13, ORANGE, True),
    ("Too slow, too manual", 18, WHITE, True),
    ("", 8, WHITE),
    ("request → weeks of planning → finally builds", 13, MUTED, False),
    ("", 8, WHITE),
    ("Missing: leverage, speed — the actual\nbenefit of AI agents.", 13, WHITE),
    ("", 8, WHITE),
    ("Result: AI becomes expensive autocomplete\nin a process designed for humans.", 13, MUTED),
])

hline(s, 5.7)
multiline_box(s, 0.6, 5.9, 12, 0.9, [
    ("One shared root:   ", 16, MUTED, False),
    ("the maker grades their own work.", 16, BLUE, True),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — THE FIX: MAKER ≠ CHECKER
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 3)
label(s, 0.6, 0.4, "02 / THE FIX")
box(s, 0.6, 0.8, 11, 0.8, "Maker ≠ Checker", size=36, bold=True)
hline(s, 1.8)

# three columns
cols = [
    (0.6,  "SPECIALIST BUILDS",
     'Owns the files.\nWrites tests for its slice.\n"Done" = tests pass,\nwith output as evidence.'),
    (4.7,  "ADVERSARY ATTACKS",
     "A separate role, rewarded\nfor finding the break —\nraces, edges, missing scope,\nsecurity holes."),
    (8.8,  "EVIDENCE GATE",
     "Passes only on proof:\ntest output, a diff,\na screenshot —\nnever a claim."),
]
for left, title, body in cols:
    box(s, left, 2.1, 3.8, 0.45, title, size=12, color=BLUE, bold=True)
    box(s, left, 2.6, 3.8, 2.5, body, size=15, color=WHITE)

hline(s, 5.3)
multiline_box(s, 0.6, 5.5, 12, 1.5, [
    ("Evidence over assertion", 15, ORANGE, True),
    ("test output  ·  a diff  ·  a screenshot  ·  never a claim", 14, MUTED, False),
    ("", 8, WHITE),
    ("An agent asked to both build and verify will fit the verification to the implementation.\n"
     "Self-certification is rationalization. The structural separation is the guarantee.", 13, MUTED),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — STATE MACHINE
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 4)
label(s, 0.6, 0.4, "03 / THE PROJECT LIFECYCLE")
box(s, 0.6, 0.8, 11, 0.8, "Five states. One direction.", size=28, bold=True)
hline(s, 1.8)

states = [
    ("not_initialized", "No framework yet. Run hssd init.", MUTED),
    ("initialized",     "Framework installed. Project detected.\nReview or write project.md.", WHITE),
    ("identified",      "project.md approved by human.\nVision, objectives, non-goals locked.", BLUE),
    ("architected",     "ADR approved by human.\nStack, data model, decisions locked.", BLUE),
    ("operational  ∞",  "Intake → iteration → engineering loop.\nRuns forever. Never \"done\".", GREEN),
]

top = 2.1
for name, desc, col in states:
    box(s, 0.6, top, 3.5, 0.35, name, size=14, color=col, bold=True)
    box(s, 4.2, top, 8.5, 0.35, desc, size=13, color=MUTED)
    top += 0.82

hline(s, top + 0.05)
multiline_box(s, 0.6, top + 0.2, 12, 0.8, [
    ("Smart init:  ", 13, BLUE, True),
    ("hssd init detects blank vs existing project automatically. "
     "A new codebase starts from scratch; an existing one gets a draft project.md to review and approve.", 13, MUTED),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — OPERATIONAL CYCLE
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 5)
label(s, 0.6, 0.4, "04 / THE OPERATIONAL CYCLE")
box(s, 0.6, 0.8, 11, 0.7,
    "Once operational, the project never finishes — it cycles.", size=24, bold=True)
hline(s, 1.7)

code(s, 0.6, 1.9, 12.1, 2.0,
     "  DEMAND  ──►  INTAKE  ──►  ITERATION PLANNING  ──►  ITERATIONS  ──►  CONVERGE\n"
     "    ▲                                                                       │\n"
     "    └────────────────────── next demand ◄──────────────────────────────────┘\n"
     "                   ▲\n"
     "            JANITOR (always on) — discovers drift, debt, bugs → files new intakes")

stages = [
    ("INTAKE",
     "Any demand enters here — manual, structured\n(Linear/JSON), unstructured (Slack/email),\nor operational (janitor finding debt/bugs).\n"
     "3-stage grooming: demand → architecture lite → split."),
    ("ITERATION",
     "Scope-bounded unit. Stories run through the\nfull P0→P4 engineering loop, each in its own\nworktree. Variadic: activate 1 or 20 in parallel."),
    ("CONVERGE",
     "Verified worktrees merge back to main\nin dependency order. Janitor rescans.\nNext demand enters intake."),
]
left = 0.6
for title, body in stages:
    box(s, left, 4.1, 4.0, 0.4, title, size=12, color=ORANGE, bold=True)
    box(s, left, 4.55, 4.0, 2.1, body, size=13, color=WHITE)
    left += 4.27

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — INTAKE CYCLE
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 6)
label(s, 0.6, 0.4, "05 / THE INTAKE CYCLE")
box(s, 0.6, 0.8, 11, 0.7, "Any demand, any form. One governed path.", size=24, bold=True)
hline(s, 1.7)

# four intake forms
forms = [
    ("Manual", "A brief written by the\ndeveloper — free-form,\njust enough to start."),
    ("Structured", "Linear/Jira export, JSON,\ndesign spec — machine\nor human authored."),
    ("Unstructured", "Slack message, email,\ncall transcript — the\njanitor or AI parses it."),
    ("Operational", "Janitor finding — drift,\ndebt, latent bug — filed\nautomatically."),
]
left = 0.6
for title, body in forms:
    box(s, left, 1.9, 2.9, 0.4, title, size=13, color=BLUE, bold=True)
    box(s, left, 2.35, 2.9, 1.4, body, size=13, color=WHITE)
    left += 3.1

hline(s, 3.9)
box(s, 0.6, 4.05, 12, 0.35, "3-stage grooming cycle", size=14, color=ORANGE, bold=True)

stages2 = [
    ("1. GROOMING",
     "product-analyst decomposes the demand.\ndefinition-skeptic gates: testable? complete?\nin scope? Unresolved → assumptions logged."),
    ("2. ARCHITECTURE LITE",
     "architect checks: does this fit the ADR?\nNew decisions extend the ADR (human approves\naddendum). No new architecture by surprise."),
    ("3. SPLIT",
     "story-writer breaks the intake into stories.\nEach story enters the backlog with verifiable\nacceptance criteria. Ready for iteration planning."),
]
left = 0.6
for title, body in stages2:
    box(s, left, 4.45, 3.9, 0.4, title, size=12, color=MUTED, bold=True)
    box(s, left, 4.9, 3.9, 2.1, body, size=13, color=WHITE)
    left += 4.25

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — ITERATIONS
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 7)
label(s, 0.6, 0.4, "06 / ITERATIONS")
box(s, 0.6, 0.8, 11, 0.7,
    "Scope-bounded, variadic, parallel by default.", size=24, bold=True)
hline(s, 1.7)

code(s, 0.6, 1.9, 12.1, 1.3,
     "  # Individual contributor\n"
     "  hssd iteration activate iter-001\n\n"
     "  # Team / parallel workstreams\n"
     "  hssd iteration activate iter-001 iter-002 iter-003\n\n"
     "  # AI fleet — 20 screens, 20 parallel governance loops\n"
     "  hssd iteration activate iter-001 iter-002 ... iter-020")

multiline_box(s, 0.6, 3.4, 12, 0.8, [
    ("Parallelism is implicit in how many IDs you pass. No flags, no modes.\n"
     "The caller — Claude Code, hssd engage, or a CI pipeline — manages orchestration.", 14, MUTED),
])

hline(s, 4.3)

items = [
    ("Each iteration", "runs in its own git worktree.\nStories are fully isolated."),
    ("Governance loop", "every story runs P0→P4,\nincluding 5 adversaries."),
    ("Human gates", "Spec Lock and Merge per story.\nThe developer stays in the loop."),
    ("Convergence", "Merge in dependency order.\nJanitor rescans after each merge."),
]
left = 0.6
for title, body in items:
    box(s, left, 4.5, 2.9, 0.35, title, size=12, color=BLUE, bold=True)
    box(s, left, 4.9, 2.9, 1.4, body, size=13, color=WHITE)
    left += 3.15

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — ENGINEERING LOOP
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 8)
label(s, 0.6, 0.4, "07 / THE ENGINEERING LOOP")
box(s, 0.6, 0.8, 11, 0.7, "P0 → P4. Every story, every time. No shortcuts.", size=24, bold=True)
hline(s, 1.7)

phases = [
    ("P0", "Intake",            "Understand the story, surface assumptions, log open questions."),
    ("P1", "Stories & AC",      "Write verifiable acceptance criteria (Gherkin). These become the tests."),
    ("P2", "Architecture",      "Design the implementation — not the code. Story-level, inherits ADR."),
    ("◆",  "SPEC LOCK",         "Human approves the locked spec. No code before this gate."),
    ("P3a","Red",               "Test-author writes tests from the locked AC. They MUST fail first."),
    ("P3b","Green",             "Implement until tests pass. Loop P3b until green."),
    ("P4", "Verify",            "5 independent adversaries attack the result. Any BLOCK → back to P3b."),
    ("◆",  "MERGE",             "Human reviews the evidence — green tests + adversary verdicts — then approves."),
]

top = 1.95
for phase, name, desc in phases:
    col = ORANGE if phase == "◆" else BLUE
    bld = phase == "◆"
    box(s, 0.6,  top, 0.6,  0.35, phase, size=13, color=col, bold=True)
    box(s, 1.25, top, 2.2,  0.35, name,  size=13, color=WHITE, bold=bld)
    box(s, 3.6,  top, 9.2,  0.35, desc,  size=13, color=MUTED)
    top += 0.57

hline(s, top + 0.05)
box(s, 0.6, top + 0.2, 12, 0.5,
    "P4 is loop-until-dry: any adversary BLOCK returns to P3b. Done = green AND all adversaries pass.",
    size=13, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — P4 ADVERSARIAL FAN-OUT
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 9)
label(s, 0.6, 0.4, "08 / P4 — THE ADVERSARIAL FAN-OUT")
box(s, 0.6, 0.8, 11, 0.7,
    "loop-until-dry: iterate until the adversaries find nothing.", size=22, bold=True)
hline(s, 1.7)

adversaries = [
    ("security-adversary",      "SQL injection · prompt injection · auth abuse · secret leak"),
    ("independent-verifier",    "Runs every acceptance criterion against every test. Evidence required."),
    ("completion-challenger",   "Proves it is NOT done — cut scope, stubs, TODOs, happy path only."),
    ("test-adversary",          "Fires real concurrency races. Hunts tests that pass for the wrong reason."),
    ("regression-hunter",       "One question: what does this change break?"),
]

top = 2.0
for name, desc in adversaries:
    box(s, 0.6, top, 4.0, 0.35, name, size=13, color=ORANGE, bold=True)
    box(s, 4.8, top, 8.0, 0.35, desc, size=13, color=WHITE)
    top += 0.82

hline(s, top + 0.1)
multiline_box(s, 0.6, top + 0.3, 12, 1.0, [
    ("PASS  ", 15, GREEN, True),
    ("=  all acceptance criteria green  AND  the adversaries find nothing.", 15, WHITE, False),
    ("", 8, WHITE),
    ("The maker never declares victory. The independent checkers do — and only on evidence.", 13, MUTED),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — 14 ROLES
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 10)
label(s, 0.6, 0.4, "09 / THE ROLE SYSTEM")
box(s, 0.6, 0.8, 11, 0.7, "14 specialized roles, each in isolation.", size=24, bold=True)
hline(s, 1.7)

col1 = [
    ("INTAKE", BLUE),
    ("product-analyst", "Decomposes demand into concerns"),
    ("definition-skeptic", "Validates testability, completeness, scope"),
    ("", None),
    ("ARCHITECTURE", BLUE),
    ("architect", "Proposes design, justifies decisions"),
    ("architecture-adversary", "Finds failure modes, attacks the design"),
    ("", None),
    ("DELIVERY", BLUE),
    ("story-writer", "Writes stories with verifiable AC"),
    ("ac-adversary", "Attacks the acceptance criteria"),
    ("test-author", "Writes failing tests from the spec"),
    ("backend-dev / frontend-dev", "Implements until green"),
]

col2 = [
    ("P4 VERIFICATION", BLUE),
    ("independent-verifier", "Every AC ↔ every test"),
    ("completion-challenger", "Proves NOT done"),
    ("test-adversary", "Vacuous tests, races, false confidence"),
    ("regression-hunter", "What broke?"),
    ("security-adversary", "Injection, auth, secret leaks"),
    ("", None),
    ("CONTINUOUS", BLUE),
    ("janitor", "Drift, debt, latent bugs → new intakes"),
]

top = 1.95
for name, desc in col1:
    if desc is None:
        top += 0.18
        continue
    if isinstance(desc, RGBColor):
        box(s, 0.6, top, 6.0, 0.32, name, size=11, color=desc, bold=True)
    else:
        box(s, 0.6, top, 2.8, 0.32, name, size=12, color=WHITE)
        box(s, 3.5, top, 3.2, 0.32, desc, size=11, color=MUTED)
    top += 0.38

top = 1.95
for name, desc in col2:
    if desc is None:
        top += 0.18
        continue
    if isinstance(desc, RGBColor):
        box(s, 7.0, top, 6.0, 0.32, name, size=11, color=desc, bold=True)
    else:
        box(s, 7.0, top, 2.5, 0.32, name, size=12, color=WHITE)
        box(s, 9.6, top, 3.6, 0.32, desc, size=11, color=MUTED)
    top += 0.38

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — EVIDENCE OVER ASSERTION
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 11)
label(s, 0.6, 0.4, "10 / EVIDENCE OVER ASSERTION")
box(s, 0.6, 0.8, 11, 0.7,
    '"Done" is never a claim. It is always evidence.', size=24, bold=True)
hline(s, 1.7)

rows = [
    ('What "done" means',          "Evidence"),
    ("Tests were real",            "Red log — they failed before implementation"),
    ("Implementation works",       "Green log — actual test output"),
    ("Every AC is covered",        "independent-verifier verdict"),
    ("Nothing was missed",         "completion-challenger verdict"),
    ("Tests are not vacuous",      "test-adversary verdict"),
    ("Nothing broke",              "regression-hunter verdict"),
    ("Code is safe",               "security-adversary verdict"),
]

top = 1.95
for i, (what, evidence) in enumerate(rows):
    is_header = i == 0
    col_w = WHITE if is_header else MUTED
    col_e = BLUE  if is_header else WHITE
    sz    = 13    if is_header else 13
    bld   = is_header
    box(s, 0.6, top, 5.5, 0.38, what,     size=sz, color=col_w, bold=bld)
    box(s, 6.3, top, 6.8, 0.38, evidence, size=sz, color=col_e, bold=bld)
    top += 0.5

hline(s, top + 0.1)
box(s, 0.6, top + 0.25, 12, 0.4,
    "Evidence is written to .harness/engagements/<id>/ during the loop — not reconstructed at the end.",
    size=12, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — NEW ROLE OF THE DEVELOPER
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 12)
label(s, 0.6, 0.4, "11 / THE NEW ROLE")
box(s, 0.6, 0.8, 11, 0.7,
    "You stop writing code. You run the system that writes it.", size=22, bold=True)
hline(s, 1.7)

multiline_box(s, 0.6, 1.9, 5.7, 5.0, [
    ("What you stop doing", 14, ORANGE, True),
    ("", 6, WHITE),
    ("Writing first drafts of code", 13, MUTED),
    ("Writing first drafts of tests", 13, MUTED),
    ("Manual code review for logic errors\n(adversaries do this)", 13, MUTED),
    ("Wondering if you missed something\n(P4 tells you)", 13, MUTED),
    ("", 8, WHITE),
    ("What you become", 14, ORANGE, True),
    ("", 6, WHITE),
    ("The person who defines objectives clearly\nenough for an agent to act on them", 13, WHITE),
    ("The person who approves spec at Spec Lock", 13, WHITE),
    ("The person who reviews evidence at Merge", 13, WHITE),
    ("The person who decides the next intake", 13, WHITE),
])

multiline_box(s, 7.0, 1.9, 5.8, 5.0, [
    ("What does not change", 14, BLUE, True),
    ("", 6, WHITE),
    ("Accountability.\nThe developer's signature is on every merge.", 13, WHITE),
    ("", 8, WHITE),
    ("Judgment.\nThe human gates exist because some\ndecisions require human judgment,\nnot just correctness.", 13, WHITE),
    ("", 8, WHITE),
    ("Architecture.\nThe ADR is yours. The agents advise.\nYou decide.", 13, WHITE),
    ("", 8, WHITE),
    ("Build the loop. Stay the engineer.", 14, BLUE, True),
])

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — GETTING STARTED
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 13)
label(s, 0.6, 0.4, "12 / GETTING STARTED")
box(s, 0.6, 0.8, 11, 0.7, "Zero to operational in four commands.", size=24, bold=True)
hline(s, 1.7)

code(s, 0.6, 1.9, 5.7, 2.6,
     "# New project\npip install harness-studio\nhssd init\n"
     "# → write .harness/project.md\nhssd project approve\n"
     "# → architecture review\nhssd architecture approve\n"
     "# → first intake\nhssd intake add brief.md\nhssd iteration activate <id>")

code(s, 6.9, 1.9, 5.8, 2.6,
     "# Existing project\nhssd init\n"
     "# → review AI-drafted project.md\nhssd project approve\n"
     "# → review AI-drafted ADR\nhssd architecture approve\n"
     "# → first intake (often:\n#    address architecture findings)\nhssd intake add brief.md\nhssd iteration activate <id>")

hline(s, 4.7)
box(s, 0.6, 4.9, 12, 0.35, "Then, continuously:", size=13, color=ORANGE, bold=True)
code(s, 0.6, 5.3, 12.1, 1.4,
     "  hssd intake add <any demand — manual, structured, unstructured, or janitor-filed>\n"
     "  hssd iteration plan --stories <ids>        # group stories into iterations\n"
     "  hssd iteration activate id1 id2 id3        # start N parallel engineering loops\n"
     "  hssd status                                 # see everything in flight")

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — THREE LAYERS OF ROLE QUALITY
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 14, total=15)
label(s, 0.6, 0.4, "13 / SKILL COMPOSITION")
box(s, 0.6, 0.8, 11, 0.7, "Three layers of role quality.", size=24, bold=True)
hline(s, 1.7)

layers = [
    ("LAYER 1", "Agent card", ".claude/agents/<name>.md",
     "Defines WHAT the role IS — identity, scope, output type, reward signal."),
    ("LAYER 2", "Role skill", ".claude/skills/roles/<name>/SKILL.md",
     "Defines HOW the role EXECUTES — non-negotiables, exact output format, failure modes, loop discipline."),
    ("LAYER 3", "Engineering skills", ".claude/skills/<name>/SKILL.md",
     "Defines WHAT the role must CATCH — domain knowledge: concurrency, indexing, UTC, APIs, complexity."),
]
top = 1.9
for tag, title, path, desc in layers:
    box(s, 0.6, top, 1.1, 0.32, tag, size=11, color=ORANGE, bold=True)
    box(s, 1.8, top, 2.5, 0.32, title, size=13, color=WHITE, bold=True)
    box(s, 4.4, top, 4.2, 0.32, path, size=10, color=GREEN)
    box(s, 0.6, top + 0.35, 12.6, 0.32, desc, size=12, color=MUTED)
    top += 0.95

hline(s, top + 0.05)
box(s, 0.6, top + 0.2, 12.6, 0.32,
    "Composition rule: every subagent prompt = role skill + relevant engineering skills + task description.",
    size=13, color=BLUE, bold=False)

top2 = top + 0.65
rows = [
    ("architect / architecture-adversary", "complexity-guard · sqlite-concurrency · sql-indexing · datetime-utc · api-conventions · api-design · resilience · push-over-pull"),
    ("backend-dev",                        "python · fastapi · sqlite-concurrency · sql-indexing · datetime-utc · api-conventions · resilience · complexity-guard"),
    ("test-author",                        "python / typescript · sqlite-concurrency · api-conventions"),
    ("security-adversary",                 "api-conventions · api-design"),
    ("janitor",                            "ALL engineering skills — scans for violations of every convention"),
    ("intake / AC / completion roles",     "(none — problem framing and completeness checking, not technical domain)"),
]
for role, skills in rows:
    box(s, 0.6,  top2, 3.8, 0.32, role,   size=11, color=ORANGE, bold=False)
    box(s, 4.55, top2, 8.8, 0.32, skills, size=11, color=MUTED,  bold=False)
    top2 += 0.42

# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — CLOSING / PRINCIPLES
# ─────────────────────────────────────────────────────────────────────────────
s = new_slide()
slide_number(s, 15, total=15)
box(s, 0.6, 0.5, 12, 0.5, "// THE BET", size=13, color=MUTED)
multiline_box(s, 0.6, 1.1, 10, 1.5, [
    ("Stop hoping the AI got it right.", 36, WHITE, True),
    ("Prove it.", 36, BLUE, True),
])
hline(s, 2.9)

principles = [
    ("Evidence over assertion.",
     "Done means test output, a diff, or a verdict — never a claim."),
    ("Maker ≠ checker.",
     "The agent that builds never certifies its own work."),
    ("Governance is standing, not an end task.",
     "Every intake, every story, every merge — the loop is always on."),
    ("Simple scales.",
     "hssd iteration activate id1 id2 id3 starts 3 parallel processes. No flags."),
    ("Documentation is the product.",
     "A tool without clear docs is a tool no one uses. Rigor means nothing if no one can find it."),
]

top = 3.1
for title, body in principles:
    box(s, 0.6, top, 5.5, 0.32, title, size=13, color=ORANGE, bold=True)
    box(s, 6.3, top, 6.8, 0.32, body,  size=13, color=MUTED)
    top += 0.62

hline(s, top + 0.1)
box(s, 0.6, top + 0.25, 8, 0.4,
    "github.com/harness-studio/harness-studio   ·   hssd   ·   v0.2.0-alpha.1",
    size=11, color=MUTED)

# ─────────────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "Harness-Studio-Overview.pptx")
prs.save(out)
print(f"Saved: {out}")
